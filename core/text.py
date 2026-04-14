from __future__ import annotations

import io
import logging
from typing import Dict, Optional

import re
from ingestion.extraction.text_cleaning import clean_extracted_text, clean_markdown_text, merge_native_and_ocr_text, text_looks_weak
from src.clients.llm import JDPChatOpenAI

logger = logging.getLogger(__name__)

try:
    from markitdown import MarkItDown  # type: ignore
except Exception: # pragma: no cover - repo/runtime dependent
    MarkItDown = None  # type: ignore


def extract_markdown(
    file_bytes: bytes,
    filename: str,
    *,
    enable_ocr: bool = True,
    ocr_page_limit: int = 5,
    ocr_image_limit: int = 8,
) -> str:
    """
    Convert file bytes to Markdown via MarkItDown, clean it, and optionally OCR weak image-heavy files.

    Drop-in replacement for the current helper.
    """
    if MarkItDown is None:
        raise ImportError("markitdown is required: pip install 'markitdown[all]'")

    md = MarkItDown()
    stream = io.BytesIO(file_bytes)
    stream.name = filename

    raw_text = ""
    try:
        raw_text = md.convert_stream(stream).text_content or ""
    except Exception as exc:
        logger.warning("MarkItDown extraction failed for %s: %s", filename, exc)

    cleaned = clean_markdown_text(raw_text)
    if not enable_ocr:
        return cleaned

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("pdf", "pptx", "ppt"):
        return cleaned

    if not text_looks_weak(cleaned):
        return cleaned

    ocr_text = _try_ocr_fallback(
        file_bytes,
        filename,
        page_limit=ocr_page_limit,
        image_limit=ocr_image_limit,
    )
    merged = merge_native_and_ocr_text(cleaned, ocr_text)
    if merged != cleaned:
        logger.info("OCR fallback added text for %s", filename)
    return merged


def word_count(text: str) -> int:
    return len((text or "").split())


def _try_ocr_fallback(
    file_bytes: bytes,
    filename: str,
    *,
    page_limit: int = 5,
    image_limit: int = 8,
) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        if ext == "pdf":
            return ocr_pdf_bytes(file_bytes, page_limit=page_limit)
        if ext == "pptx":
            return ocr_pptx_bytes(file_bytes, image_limit=image_limit)
    except Exception as exc: # pragma: no cover - best effort path
        logger.warning("OCR fallback failed for %s: %s", filename, exc)
    return ""


def ocr_pdf_bytes(file_bytes: bytes, *, page_limit: int = 5) -> str:
    try:
        import fitz  # PyMuPDF
        import pytesseract
    except Exception:
        return ""

    parts: list[str] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page_num in range(min(len(doc), page_limit)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_bytes = pix.tobytes("png")
            text = pytesseract.image_to_string(io.BytesIO(img_bytes))
            if text and text.strip():
                parts.append(text)
    return clean_markdown_text("\n\n".join(parts))


def ocr_pptx_bytes(file_bytes: bytes, *, image_limit: int = 8) -> str:
    try:
        import pytesseract
        from PIL import Image
        from pptx import Presentation
    except Exception:
        return ""

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    seen = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            image = getattr(shape, "image", None)
            if image is None:
                continue
            try:
                img = Image.open(io.BytesIO(image.blob))
                text = pytesseract.image_to_string(img)
                if text and text.strip():
                    parts.append(text)
            except Exception:
                continue
            seen += 1
            if seen >= image_limit:
                return clean_markdown_text("\n\n".join(parts))
    return clean_markdown_text("\n\n".join(parts))


def clean_opt_text(raw_text: str) -> str:
    return clean_extracted_text(raw_text)


_SECTION_NAMES = [
    "Idea Card Executive Summary",
    "Problem Statement/Market Opportunity",
    "Business Solution and Objectives",
    "Alternative Solutions",
    "Value Proposition & Key Metrics",
    "Interdependencies",
    "Estimated Costs",
    "Resources/Investments Needed for Business Case",
]


def extract_signal_sections(raw_text: str) -> Dict[str, str]:
    cleaned = clean_opt_text(raw_text)
    lower = cleaned.lower()

    positions = []
    for name in _SECTION_NAMES:
        idx = lower.find(name.lower())
        if idx != -1:
            positions.append((idx, name))

    positions.sort()
    sections: Dict[str, str] = {}
    for i, (start, name) in enumerate(positions):
        end = positions[i + 1][0] if i + 1 < len(positions) else len(cleaned)
        sections[name] = cleaned[start:end].strip()
    return sections


def normalize_for_search(text: Optional[str], max_chars: int = 2500) -> str:
    cleaned = clean_opt_text(text)
    cleaned = cleaned.lower()
    cleaned = re.sub(r"[^a-z0-9\n\.\,\/\-\&\$\+\ ]+", " ", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned[:max_chars]


import re
from typing import Any, Dict, List
from .prompts import safe_json_extract


SIGNAL_TERMS = {
    "member", "provider", "employer", "patient", "care", "clinical",
    "product", "offering", "launch", "enroll", "enrollment", "billing",
    "invoice", "payment", "quote", "pricing", "network", "program",
    "engagement", "support", "workflow", "digital", "platform", "risk",
    "compliance", "authorization", "utilization", "claims", "appeal",
    "navigation", "coordination", "sales", "growth", "commercial"
}

ACTION_TERMS = {
    "launch", "improve", "reduce", "increase", "streamline", "enable",
    "automate", "support", "expand", "build", "create", "optimize",
    "establish", "manage", "onboard", "configure", "fulfill", "issue",
    "resolve", "align", "execute"
}


def extract_signal_lines(text: str, max_lines: int = 60, max_chars: int = 12000) -> str:
    """
    Pull the most business-salient lines from the PPT text so the summary step
    sees dense signal instead of the whole noisy deck.
    """
    lines = []
    for raw in (text or "").splitlines():
        line = re.sub(r"\s+", " ", raw).strip()
        if not line:
            continue
        if len(line) < 4:
            continue
        lines.append(line)

    scored = []
    seen = set()

    for idx, line in enumerate(lines):
        low = line.lower()
        if low in seen:
            continue
        seen.add(low)

        score = 0

        # Headings / short statements are often high signal
        if len(line) < 80:
            score += 1

        # Bullet-ish lines
        if raw_starts_like_bullet(line):
            score += 2

        # Signal vocabulary
        score += sum(2 for term in SIGNAL_TERMS if term in low)
        score += sum(1 for term in ACTION_TERMS if term in low)

        # Numeric business lines often matter
        if re.search(r"\$\d+|\d+ ?%", line):
            score += 1

        # Penalize extremely long / likely noisy paragraphs
        if len(line) > 240:
            score -= 2

        if score > 0:
            scored.append((score, idx, line))

    # Keep strongest lines, but preserve original order
    top = sorted(scored, key=lambda x: (x[0], x[1]))[-max_lines:]
    top = sorted(top, key=lambda x: x[1])

    out_lines = []
    total = 0
    for _, _, line in top:
        if total + len(line) + 1 > max_chars:
            break
        out_lines.append(line)
        total += len(line) + 1

    return "\n".join(out_lines)


def raw_starts_like_bullet(line: str) -> bool:
    return bool(re.match(r"^[-*\•\d\)]", line))


def normalize_ppt_summary(parsed: Any, fallback_text: str) -> Dict[str, Any]:
    if not isinstance(parsed, dict):
        parsed = {}

    def _as_list(value) -> List[str]:
        if value is None:
            return []
        if isinstance(value, list):
            return [str(x).strip() for x in value if str(x).strip()]
        parts = [p.strip() for p in re.split(r"[,;\n]+", value) if p.strip()]
        return parts

    summary = {
        "short_summary": str(parsed.get("short_summary") or "").strip(),
        "business_goal": str(parsed.get("business_goal") or "").strip(),
        "actors": _as_list(parsed.get("actors")),
        "direct_functions": _as_list(parsed.get("direct_functions")),
        "implied_functions": _as_list(parsed.get("implied_functions")),
        "change_types": _as_list(parsed.get("change_types")),
        "domain_tags": _as_list(parsed.get("domain_tags")),
        "evidence_sentences": _as_list(parsed.get("evidence_sentences")),
    }

    if not summary["short_summary"]:
        summary["short_summary"] = fallback_text[:300].strip()

    if not summary["business_goal"]:
        summary["business_goal"] = summary["short_summary"]

    return summary


def generate_ppt_semantic_summary(cleaned_text: str) -> Dict[str, Any]:
    """
    Create a compact semantic profile of the raw PPT for retrieval.
    """
    signal_text = extract_signal_lines(cleaned_text)

    prompt = f"""
You convert noisy enterprise healthcare PPT text into a compact semantic retrieval profile.

Return VALID JSON ONLY with exactly these keys:
short_summary
business_goal
actors
direct_functions
implied_functions
change_types
domain_tags
evidence_sentences

Rules:
- Be specific, not generic.
- actors must be a JSON array of strings.
- direct_functions must contain only functions directly stated or strongly evidenced.
- implied_functions may include downstream enterprise processes that are clearly implied.
- evidence_sentences must be short verbatim or near-verbatim lines from the provided text.
- Keep short_summary to 1-2 sentences max.

PPT_TEXT:
{signal_text}
""".strip()

    reply = JDPChatOpenAI(model="gpt-4o-mini-idp").invoke(input=prompt)

    raw = getattr(reply, "content", "") or ""
    parsed = safe_json_extract(raw)
    summary = normalize_ppt_summary(parsed, signal_text)

    summary["_signal_text"] = signal_text
    summary["_raw_summary_response"] = raw
    return summary


def build_summary_retrieval_text(summary: Dict[str, Any]) -> str:
    parts = []

    if summary.get("short_summary"):
        parts.append(summary["short_summary"])

    if summary.get("business_goal"):
        parts.append(f"Business Goal: {summary['business_goal']}")

    if summary.get("actors"):
        parts.append("Actors: " + ", ".join(summary["actors"]))

    if summary.get("direct_functions"):
        parts.append("Direct Functions: " + ", ".join(summary["direct_functions"]))

    if summary.get("implied_functions"):
        parts.append("Implied Functions: " + ", ".join(summary["implied_functions"]))

    if summary.get("change_types"):
        parts.append("Change Types: " + ", ".join(summary["change_types"]))

    if summary.get("domain_tags"):
        parts.append("Domain Tags: " + ", ".join(summary["domain_tags"]))

    return "\n".join(parts)