from __future__ import annotations

import io
import logging
from typing import Optional

from .text_cleaning import clean_markdown_text, merge_native_and_ocr_text, text_looks_weak


logger = logging.getLogger(__name__)


try:
    from markitdown import MarkItDown  # type: ignore
except Exception:  # pragma: no cover - repo/runtime dependent
    MarkItDown = None  # type: ignore


def extract_markdown(
    file_bytes: bytes,
    filename: str,
    *,
    enable_ocr: bool = True,
    ocr_page_limit: int = 5,
    ocr_image_limit: int = 8,
) -> str:
    """
    Convert file bytes to Markdown via MarkItDown, clean it, and optionally OCR weak image-heavy files.

    Drop-in replacement for the current helper.
    """
    if MarkItDown is None:
        raise ImportError("markitdown is required: pip install 'markitdown[all]'")

    md = MarkItDown()
    stream = io.BytesIO(file_bytes)
    stream.name = filename

    raw_text = ""
    try:
        raw_text = md.convert_stream(stream).text_content or ""
    except Exception as exc:
        logger.warning("MarkItDown extraction failed for %s: %s", filename, exc)

    cleaned = clean_markdown_text(raw_text)
    if not enable_ocr:
        return cleaned

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("pdf", "pptx", "ppt"):
        return cleaned

    if not text_looks_weak(cleaned):
        return cleaned

    ocr_text = _try_ocr_fallback(
        file_bytes,
        filename,
        page_limit=ocr_page_limit,
        image_limit=ocr_image_limit,
    )
    
    merged = merge_native_and_ocr_text(cleaned, ocr_text)
    if merged != cleaned:
        logger.info("OCR fallback added text for %s", filename)
    return merged


def word_count(text: str) -> int:
    return len((text or "").split())


def _try_ocr_fallback(
    file_bytes: bytes,
    filename: str,
    *,
    page_limit: int = 5,
    image_limit: int = 8,
) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        if ext == "pdf":
            return _ocr_pdf_bytes(file_bytes, page_limit=page_limit)
        if ext == "pptx":
            return _ocr_pptx_bytes(file_bytes, image_limit=image_limit)
    except Exception as exc:  # pragma: no cover - best effort path
        logger.warning("OCR fallback failed for %s: %s", filename, exc)
    return ""


def _ocr_pdf_bytes(file_bytes: bytes, *, page_limit: int = 5) -> str:
    try:
        import fitz  # PyMuPDF
        import pytesseract
    except Exception:
        return ""

    parts: list[str] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page_num in range(min(len(doc), page_limit)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_bytes = pix.tobytes("png")
            text = pytesseract.image_to_string(io.BytesIO(img_bytes))
            if text and text.strip():
                parts.append(text)
    return clean_markdown_text("\n\n".join(parts))


def _ocr_pptx_bytes(file_bytes: bytes, *, image_limit: int = 8) -> str:
    try:
        import pytesseract
        from PIL import Image
        from pptx import Presentation
    except Exception:
        return ""

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    seen = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            image = getattr(shape, "image", None)
            if image is None:
                continue
            try:
                img = Image.open(io.BytesIO(image.blob))
                text = pytesseract.image_to_string(img)
                if text and text.strip():
                    parts.append(text)
            except Exception:
                continue
            seen += 1
            if seen >= image_limit:
                return clean_markdown_text("\n\n".join(parts))
    return clean_markdown_text("\n\n".join(parts))


def extract_image_text(file_bytes: bytes, *, filename: str = "image") -> str:
    """OCR plain image bytes and return cleaned text."""
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return ""

    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        cleaned = clean_markdown_text(text)
        if not cleaned:
            pass
        return cleaned
    except Exception as exc:  # pragma: no cover - best effort path
        logger.warning("Image OCR failed for %s: %s", filename, exc)
        return ""